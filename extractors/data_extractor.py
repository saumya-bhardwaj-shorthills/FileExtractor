import fitz  # PyMuPDF for PDF handling
import camelot  # For PDF table extraction
import docx  # For DOCX handling
import pptx  # For PPTX handling


class DataExtractor:
    def __init__(self, loader):
        self.loader = loader
        self.file_path = loader.file_path

    def extract_text(self):
        if self.file_path.endswith('.pdf'):
            # Extract text from PDF
            reader = self.loader.load_file()
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text

        elif self.file_path.endswith('.docx'):
            # Extract text from DOCX
            doc = self.loader.load_file()
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text

        elif self.file_path.endswith('.pptx'):
            # Extract text from PPTX
            ppt = self.loader.load_file()
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            raise ValueError("Unsupported file format for text extraction.")

    def extract_images(self):
        images = []

        if self.file_path.endswith('.pdf'):
            # PDF image extraction
            pdf_document = fitz.open(self.file_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                image_list = page.get_images(full=True)
                for img in image_list:
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    width, height = base_image["width"], base_image["height"]
                    images.append({
                        "image_data": image_bytes,
                        "ext": image_ext,
                        "page": page_num + 1,
                        "dimensions": (width, height)
                    })
            pdf_document.close()

        elif self.file_path.endswith('.docx'):
            # DOCX image extraction
            doc = self.loader.load_file()
            for rel in doc.part.rels:
                if "image" in doc.part.rels[rel].target_ref:
                    image = doc.part.rels[rel].target_part.blob
                    images.append(image)

        elif self.file_path.endswith('.pptx'):
            # PPTX image extraction
            ppt = self.loader.load_file()
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        images.append(shape.image)

        return images

    def extract_urls(self):
        urls = []

        if self.file_path.endswith('.pdf'):
            # PDF URL extraction
            pdf_document = fitz.open(self.file_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                links = page.get_links()
                for link in links:
                    if "uri" in link:
                        url = link["uri"]
                        rect = link["from"]
                        urls.append({
                            "url": url,
                            "page": page_num + 1,
                            "position": {
                                "x0": rect.x0,
                                "y0": rect.y0,
                                "x1": rect.x1,
                                "y1": rect.y1
                            }
                        })
            pdf_document.close()

        elif self.file_path.endswith('.docx'):
            # DOCX URL extraction
            doc = self.loader.load_file()
            for rel in doc.part.rels.values():
                if "hyperlink" in rel.target_ref:
                    urls.append(rel.target_ref)

        elif self.file_path.endswith('.pptx'):
            # URLs are rare in PPTX, so skipping this part
            pass

        return urls

    def extract_tables(self):
        tables = []

        if self.file_path.endswith('.pdf'):
            # Extract tables from PDF using Camelot
            tables = camelot.read_pdf(self.file_path, pages="all")
            table_data = [table.df for table in tables]  # List of DataFrames
            return table_data

        elif self.file_path.endswith('.docx'):
            # Extract tables from DOCX
            doc = self.loader.load_file()
            table_data = []
            for table in doc.tables:
                table_content = [[cell.text for cell in row.cells] for row in table.rows]
                table_data.append(table_content)
            return table_data

        elif self.file_path.endswith('.pptx'):
            # Extract tables from PPTX (typically tables are part of shapes)
            ppt = self.loader.load_file()
            table_data = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        table = shape.table
                        table_content = [[cell.text for cell in row.cells] for row in table.rows]
                        table_data.append(table_content)
            return table_data

        return tables
